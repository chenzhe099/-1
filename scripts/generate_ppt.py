"""
生成完整的商业计划书 PPT
- 嵌入 10 张流程图
- 自动截取 8 个模块的应用界面截图
- 专业排版，统一配色
"""
import os, sys, time, io
from docx.shared import Pt as DocxPt  # not used, for reference

# Add base path
BASE = r'd:\zhuomian\gjcx'
IMG_DIR = os.path.join(BASE, 'assets', 'flowcharts')
SCREENSHOT_DIR = os.path.join(BASE, 'assets', 'screenshots')
os.makedirs(SCREENSHOT_DIR, exist_ok=True)

APP_URL = 'http://localhost:8000'

# ====== Step 1: Take Screenshots ======
def take_screenshots():
    """Use Playwright to capture screenshots of each module"""
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print('Playwright not installed, skipping screenshots')
        return False

    print('Taking application screenshots...')
    sections = [
        ('dashboard', '数据总览驾驶舱'),
        ('disease', 'AI病虫害识别'),
        ('farming', 'AI精准农事决策'),
        ('prediction', '产量预测与农事规划'),
        ('management', '数字化农场管理'),
        ('devices', '设备监控与远程控制'),
        ('traceability', '农产品溯源管理'),
        ('permission', '权限管理与多账号协同'),
    ]

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={'width': 1440, 'height': 900})

        # First load the page
        page.goto(APP_URL, wait_until='networkidle', timeout=15000)
        time.sleep(2)

        # Navigate to each section and screenshot
        for section_id, section_name in sections:
            try:
                # Click the sidebar button
                page.click(f'.sidebar-item[data-menu="{section_id}"]')
                time.sleep(1.5)  # Wait for render + charts

                # Take full page screenshot
                path = os.path.join(SCREENSHOT_DIR, f'{section_id}.png')
                page.screenshot(path=path, full_page=False)
                print(f'  OK {section_name} -> {section_id}.png')
            except Exception as e:
                print(f'  ERR {section_name}: {e}')

        browser.close()
    return True

# ====== Step 2: Generate PPT ======
def generate_ppt():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    # Colors
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    DARK = RGBColor(44, 62, 80)
    BLUE = RGBColor(41, 128, 185)
    GRAY = RGBColor(127, 140, 141)
    LGRAY = RGBColor(189, 195, 199)
    BG = RGBColor(236, 240, 241)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK_LAYOUT = prs.slide_layouts[6]  # Blank layout

    def new_slide():
        return prs.slides.add_slide(BLANK_LAYOUT)

    def add_bg(slide, color=WHITE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.bold = bold; p.font.name = font_name
        p.alignment = align
        return txBox

    def add_title(slide, text, top=0.3):
        add_text_box(slide, 0.8, top, 11.5, 0.8, text, font_size=32, color=DARK, bold=True)

    def add_subtitle(slide, text, top=1.1):
        add_text_box(slide, 0.8, top, 11.5, 0.5, text, font_size=16, color=GRAY)

    def add_body(slide, text, left=0.8, top=1.8, width=11.5, height=5, font_size=14):
        add_text_box(slide, left, top, width, height, text, font_size=font_size, color=BLACK)

    def add_image(slide, filename, left, top, width, height=None):
        path = os.path.join(IMG_DIR, filename)
        if not os.path.exists(path):
            add_text_box(slide, left, top, width, 0.5, f'[图片: {filename}]', font_size=10, color=GRAY)
            return
        if height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))

    def add_screenshot(slide, filename, left, top, width, height=None):
        path = os.path.join(SCREENSHOT_DIR, filename)
        if not os.path.exists(path):
            add_text_box(slide, left, top, width, 0.4, f'[程序截图: {filename}]', font_size=10, color=GRAY)
            return
        if height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))

    def add_bullet_slide(slide, items, left=0.8, top=1.8, font_size=15):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11.5), Inches(5))
        tf = txBox.text_frame; tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = BLACK
            p.font.name = 'Microsoft YaHei'
            p.space_after = Pt(8)
            p.level = 0
        return txBox

    def add_section_header(slide, number, title, subtitle=''):
        """Clean section divider slide"""
        add_bg(slide, DARK)
        # Number circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(2.2), Inches(1.6), Inches(1.6))
        shape.fill.solid(); shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        tf = shape.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.text = str(number); p.font.size = Pt(48); p.font.color.rgb = WHITE
        p.font.bold = True; p.font.name = 'Arial'; p.alignment = PP_ALIGN.CENTER
        add_text_box(slide, 1, 4.2, 11.3, 1, title, font_size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            add_text_box(slide, 1, 5.2, 11.3, 0.6, subtitle, font_size=16, color=LGRAY, align=PP_ALIGN.CENTER)

    # ========================
    # SLIDE 1: COVER
    # ========================
    slide = new_slide()
    add_bg(slide, DARK)
    add_text_box(slide, 1, 1.5, 11.3, 1.2, '智 农 云', font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.8, 11.3, 0.8, 'SmartAgri Cloud', font_size=20, color=LGRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.0, 11.3, 1, '智慧农业管理系统  商业计划书', font_size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Separator line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(5.2), Inches(3.3), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE; shape.line.fill.background()
    add_text_box(slide, 1, 5.5, 11.3, 0.5, 'V1.0  |  2024年1月  |  天使轮融资  |  商业机密', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6.6, 11.3, 0.4, '智农云科技（筹）', font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ========================
    # SLIDE 2: AGENDA
    # ========================
    slide = new_slide(); add_bg(slide)
    add_title(slide, '汇报大纲')
    agenda = [
        '第一部分    项目概述与市场机会',
        '第二部分    产品演示 — 八大功能模块',
        '第三部分    技术架构与竞争优势',
        '第四部分    商业模式与收入预测',
        '第五部分    市场推广策略',
        '第六部分    团队与发展路线图',
        '第七部分    财务预测与融资需求',
    ]
    add_bullet_slide(slide, agenda, top=1.8, font_size=18)

    # ========================
    # SLIDE 3-4: SECTION 1 HEADER + OVERVIEW
    # ========================
    slide = new_slide()
    add_section_header(slide, '01', '项目概述与市场机会', '智慧农业 — 2000亿蓝海市场')

    slide = new_slide(); add_bg(slide)
    add_title(slide, '项目定位：智农云 SmartAgri Cloud')
    overview_text = [
        '核心定位：面向中大型农业企业的 SaaS 化智慧农业管理平台',
        '',
        '通过 AI（人工智能）+ IoT（物联网）技术实现农业生产全流程的数字化、智能化管理。',
        '',
        '目标市场：中国智慧农业市场预计 2025 年达到 2,000 亿元，年复合增长率超过 18%。',
        '',
        '当前痛点：农业企业数字化水平低 — 生产靠经验、病虫害发现滞后、水肥药浪费严重、',
        '食品安全追溯体系缺失、设备管理混乱。',
    ]
    add_bullet_slide(slide, overview_text, font_size=16)

    slide = new_slide(); add_bg(slide)
    add_title(slide, '六大核心痛点与解决方案')
    pain_points = [
        '痛点 1：农业生产缺乏数据支撑  →  环境传感器 + 土壤监测 + 实时数据看板',
        '痛点 2：病虫害发现滞后损失巨大  →  AI 图像识别，秒级诊断，精准防治方案',
        '痛点 3：水肥药浪费严重成本高   →  AI 精准农事决策，节水 32%，节肥 20%+',
        '痛点 4：产量不可预测靠天吃饭   →  机器学习模型预测产量，准确率 94%',
        '痛点 5：食品安全信任危机       →  全链条溯源，一物一码，区块链存证',
        '痛点 6：设备管理混乱故障频发   →  IoT 统一监控平台，预测性维护',
    ]
    add_bullet_slide(slide, pain_points, font_size=15)

    # Market size slide
    slide = new_slide(); add_bg(slide)
    add_title(slide, '市场规模与增长趋势')
    add_text_box(slide, 0.8, 1.5, 11.5, 2, '中国智慧农业市场正经历高速增长，2023 年市场规模约 1,200 亿元，预计到 2025 年将突破 2,000 亿元，2028 年达到 3,500 亿元。农业 SaaS 细分市场同期将从 80 亿元增长至 300 亿元。', font_size=14, color=BLACK)
    add_image(slide, '08_finance.png', 0.5, 3.0, 12.3, 4)

    # ========================
    # SLIDES 5-12: SECTION 2 — PRODUCT DEMO
    # ========================
    slide = new_slide()
    add_section_header(slide, '02', '产品演示 — 八大功能模块', 'AI + IoT 赋能智慧农业全场景')

    # Overview slide
    slide = new_slide(); add_bg(slide)
    add_title(slide, '产品功能全景图')
    add_image(slide, '01_arch.png', 0.3, 1.5, 12.5, 5.5)

    # Module slides — each with screenshot + description
    modules = [
        ('dashboard', '模块一：数据总览与驾驶舱',
         '实时气象监测、土壤墒情数据、设备在线状态、产量预估看板',
         '农场全貌一目了然，支持多维度数据下钻。预警中心分级推送病虫害、干旱、设备故障和极端天气四类预警，实时通知到人。'),
        ('disease', '模块二：AI 病虫害智能识别与预警',
         '手机拍照上传、AI 图像诊断、防治方案自动生成',
         '支持 50+ 常见病虫害识别，准确率超过 92%。系统自动匹配化学防治、生物防治和农业防治三重建议方案，附带安全间隔期提醒。'),
        ('farming', '模块三：AI 精准农事决策',
         '智能灌溉方案、精准施肥配方、作业工单自动生成',
         '基于土壤 NPK 养分含量和有机质水平的施肥模型，结合土壤湿度传感器的灌溉策略。节水率实测 30-35%，肥料利用率提升至 85% 以上。'),
        ('prediction', '模块四：产量预测与农事规划',
         '机器学习产量预测、智能农事排期、风险预警提示',
         '综合历史产量、气象预报、土壤数据和作物生长阶段，预测准确率 94%。提前 1-3 个月预判产量，辅助采收计划和市场销售决策。'),
        ('management', '模块五：数字化农场管理',
         '电子种植台账、人员管理、农资库存、投入产出报表',
         '覆盖播种、施肥、灌溉、采收全周期的电子化记录。多维度报表自动生成投入产出分析、人工效率统计和设备利用率报告。'),
        ('devices', '模块六：设备监控与远程控制',
         'IoT 设备统一接入、实时状态监控、远程启停控制',
         '支持灌溉泵、施肥机、传感器、控制器等主流 IoT 设备接入。基于运行时长的预测性维护算法自动生成保养工单。'),
        ('traceability', '模块七：农产品溯源管理',
         '生产全过程追溯、一物一码 QR Code、质量认证管理',
         '消费者扫码即可查看农产品从播种到采收的完整生产链路。支持农残检测、有机认证和重金属检测等质量认证信息的一站式管理。'),
        ('permission', '模块八：权限管理与多账号协同',
         'RBAC 三级权限体系、操作日志审计、多农场统一管理',
         '管理员—技术员—农户三级角色权限矩阵，支持精细化到模块的查看和编辑权限控制。完整操作日志审计满足合规要求。'),
    ]

    screenshot_map = {
        'dashboard': '01_dashboard.png', 'disease': '02_disease.png',
        'farming': '03_farming.png', 'prediction': '04_prediction.png',
        'management': '05_management.png', 'devices': '06_devices.png',
        'traceability': '07_traceability.png', 'permission': '08_permission.png'
    }
    for section_id, title, features, description in modules:
        slide = new_slide(); add_bg(slide)
        add_title(slide, title)
        # Left: screenshot
        sc_file = screenshot_map.get(section_id, f'{section_id}.png')
        add_screenshot(slide, sc_file, 0.3, 1.5, 6.8, 4.2)
        # Right: description
        add_text_box(slide, 7.5, 1.5, 5.5, 1.5, f'功能亮点：\n{features}', font_size=13, color=BLUE)
        add_text_box(slide, 7.5, 3.5, 5.5, 3.5, description, font_size=13, color=BLACK)

    # ========================
    # SLIDES 13-15: SECTION 3 — TECH & COMPETITION
    # ========================
    slide = new_slide()
    add_section_header(slide, '03', '技术架构与竞争优势', '技术壁垒 + 全模块覆盖')

    slide = new_slide(); add_bg(slide)
    add_title(slide, '技术架构层次图')
    add_image(slide, '07_tech.png', 0.3, 1.5, 12.5, 5.5)

    slide = new_slide(); add_bg(slide)
    add_title(slide, '竞争态势分析')
    comp_items = [
        '互联网巨头（阿里数农、京东农场）：品牌强但重平台轻垂直，行业响应慢',
        '传统信息化厂商（中化MAP、大北农）：行业经验丰富但技术陈旧，用户体验差',
        '硬件厂商（大疆农业、极飞）：硬件能力强，但软件平台能力薄弱',
        '创业公司（丰疆智能、麦飞科技）：灵活创新，但功能覆盖不全，融资难度大',
        '',
        '智农云的差异化优势：8 大模块全场景覆盖 + 多场景 AI 驱动 + 灵活部署方式',
        '+ 全链条溯源 + 按需定价 + 深度适配中国农情的知识库体系',
    ]
    add_bullet_slide(slide, comp_items, font_size=14)
    add_image(slide, '06_compete.png', 0.5, 5.0, 7, 2.2)

    # ========================
    # SLIDES 16-18: SECTION 4 — BUSINESS MODEL
    # ========================
    slide = new_slide()
    add_section_header(slide, '04', '商业模式与收入预测', 'SaaS 订阅 + 增值服务 + 硬件销售')

    slide = new_slide(); add_bg(slide)
    add_title(slide, '收入来源结构')
    add_image(slide, '05_bizmodel.png', 0.3, 1.5, 12.5, 5.5)

    slide = new_slide(); add_bg(slide)
    add_title(slide, 'SaaS 订阅定价策略')
    pricing = [
        '基础版 — 小型农场（<100亩）：¥9,800/年',
        '    功能：数据看板 + 农场管理 + 产品溯源',
        '',
        '专业版 — 中型农场（100-500亩）：¥48,000/年',
        '    功能：基础版全部 + 病虫害AI识别 + 农事AI决策 + 设备监控',
        '',
        '企业版 — 大型集团（500+亩）：¥120,000 ~ 500,000/年',
        '    功能：全部 8 大模块 + 私有化部署 + 定制开发 + 优先技术支持',
    ]
    add_bullet_slide(slide, pricing, font_size=15)

    # ========================
    # SLIDES 19-20: SECTION 5 — MARKETING
    # ========================
    slide = new_slide()
    add_section_header(slide, '05', '市场推广策略', '三大获客通路 + 分阶段推进')

    slide = new_slide(); add_bg(slide)
    add_title(slide, '客户获取渠道体系')
    add_image(slide, '02_channels.png', 0.3, 1.2, 12.5, 5.8)

    slide = new_slide(); add_bg(slide)
    add_title(slide, '分阶段市场策略')
    strategy = [
        '种子期（0-6 个月）：3-5 家标杆客户免费试用，打磨产品，验证 PMF',
        '    策略：联合科研机构发布白皮书，打造示范农场参观基地',
        '',
        '成长期（6-18 个月）：直销团队 + 行业展会 + 协会推荐，50-80 家付费客户',
        '    策略：农业展会参展，农资经销商渠道合作，政府智慧农业项目承接',
        '',
        '扩张期（18-36 个月）：渠道代理 + 线上投放 + 区域扩张',
        '    策略：覆盖 5+ 核心农业省份，建设 20+ 渠道代理商网络',
    ]
    add_bullet_slide(slide, strategy, font_size=14)

    # ========================
    # SLIDES 21-22: SECTION 6 — TEAM & ROADMAP
    # ========================
    slide = new_slide()
    add_section_header(slide, '06', '团队与运营规划', '13 人核心团队 + 四阶段发展路线')

    slide = new_slide(); add_bg(slide)
    add_title(slide, '核心团队组织架构（计划）')
    add_image(slide, '03_org.png', 0.3, 1.5, 12.5, 5.5)

    slide = new_slide(); add_bg(slide)
    add_title(slide, '发展路线图（2024-2027）')
    add_image(slide, '04_roadmap.png', 0.3, 1.2, 12.5, 6)

    # ========================
    # SLIDES 23-24: SECTION 7 — FINANCIALS
    # ========================
    slide = new_slide()
    add_section_header(slide, '07', '财务预测与融资需求', '3 年盈利路径 — 天使轮融资')

    slide = new_slide(); add_bg(slide)
    add_title(slide, '3 年财务预测')
    add_image(slide, '08_finance.png', 0.3, 1.2, 12.5, 5.8)

    slide = new_slide(); add_bg(slide)
    add_title(slide, '融资计划')
    fund = [
        '融资轮次：天使轮 / Pre-A',
        '融资金额：¥500 - 800 万元',
        '出让股权：10% - 15%',
        '投前估值：¥3,000 - 5,000 万元',
        '资金用途：产品研发 50% | 团队建设 30% | 市场推广 20%',
        '预计 Runway：18 个月',
        '',
        '投资亮点：',
        '    1. 2000 亿智慧农业市场，年复合增长 18%+',
        '    2. 8 大模块一站式覆盖，竞品通常仅 3-5 个',
        '    3. AI 驱动技术壁垒：图像识别 + 智能决策 + 产量预测',
        '    4. MVP 原型已完成，110+ 交互功能全部可运行演示',
        '    5. 商业模式清晰：SaaS 订阅 + 硬件 + 增值服务，多收入来源',
        '    6. 退出路径明确：农业科技并购活跃（大北农、中化、阿里均有案例）',
    ]
    add_bullet_slide(slide, fund, font_size=14)

    # ========================
    # SLIDE 25: RISK ASSESSMENT
    # ========================
    slide = new_slide(); add_bg(slide)
    add_title(slide, '风险评估与应对策略')
    add_image(slide, '09_risk.png', 0.2, 1.3, 6.5, 5.5)
    risk_items = [
        '高风险：产品落地风险',
        '  应对：与农业专家合作标注数据，迭代训练模型',
        '',
        '高风险：现金流风险',
        '  应对：严控费用 + 政府补贴 + 预收年费模式',
        '',
        '中风险：市场竞争',
        '  应对：聚焦垂直场景深度，建立行业 Know-how 壁垒',
        '',
        '中风险：客户教育成本',
        '  应对：标杆案例 + 现场演示 + 政府项目背书',
    ]
    add_bullet_slide(slide, risk_items, left=7.2, top=1.5, font_size=12)

    # ========================
    # SLIDE 26: VALUE LOOP
    # ========================
    slide = new_slide(); add_bg(slide)
    add_title(slide, '客户价值闭环体系')
    add_image(slide, '10_loop.png', 2.5, 1.5, 8.5, 5.5)

    # ========================
    # SLIDE 27: CLOSING
    # ========================
    slide = new_slide()
    add_bg(slide, DARK)
    add_text_box(slide, 1, 2.0, 11.3, 1.2, '感谢聆听', font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.5, 11.3, 0.8, '智农云 — 让每一寸土地更智慧', font_size=22, color=LGRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5.0, 11.3, 1, '期待与您深入交流', font_size=18, color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6.3, 11.3, 0.6, '[ 联系方式 ]  |  [ 产品演示预约 ]  |  [ 商业计划书完整版 ]', font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # ====== Save ======
    out_path = os.path.join(BASE, '智农云_商业计划书_演示稿.pptx')
    prs.save(out_path)
    print(f'\nPPT saved: {out_path}')
    print(f'Size: {os.path.getsize(out_path)/1024:.1f} KB')
    print(f'Slides: {len(prs.slides)}')

# ====== Main ======
if __name__ == '__main__':
    print('=' * 50)
    print('Generating Business Plan PPT...')
    print('=' * 50)

    # Step 1: Screenshots
    screenshots_ok = take_screenshots()
    if not screenshots_ok:
        print('\nScreenshots skipped (run manually if needed)')

    # Step 2: PPT
    generate_ppt()
    print('\nDone!')
